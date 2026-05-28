"""Multi-format ingestion.

Every parser produces the same `list[PageContent]` shape so the rest of
the pipeline doesn't care what the source was. Format dispatch is by
file extension; pass `format=` to override.

Supported formats (v0.1):

| Extension | Parser | Pagination |
|---|---|---|
| ``.pdf`` | pymupdf4llm | Real pages |
| ``.docx`` | python-docx | One page per H1, or whole doc if no H1 |
| ``.pptx`` | python-pptx | One page per slide |
| ``.md``, ``.markdown`` | built-in | One page per top-level H1 |
| ``.html``, ``.htm`` | markdownify | One page (post-conversion to MD) |
| ``.txt`` | built-in | One page |
| ``.png``, ``.jpg``, ``.jpeg``, ``.tif``, ``.tiff``, ``.webp`` | pymupdf | One page, `needs_ocr=True` |
"""

from __future__ import annotations

import re
from collections.abc import Callable
from pathlib import Path

import pymupdf

from kglite_docs.ingest.parser import PageContent, parse_pdf

_IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".tif", ".tiff", ".webp", ".bmp"}


def detect_format(path: str | Path) -> str:
    """Return a normalised format key for a file path (the lowercased
    extension, without the dot)."""
    return Path(path).suffix.lower().lstrip(".")


def parse_document(path: str | Path, *, format: str | None = None) -> list[PageContent]:
    """Dispatch to the right format-specific parser."""
    from kglite_docs.errors import UnsupportedFormatError
    path = Path(path)
    fmt = (format or detect_format(path)).lower()
    parser = _PARSERS.get(fmt)
    if parser is None:
        raise UnsupportedFormatError(
            f"unsupported format: {fmt!r} (file {path}). "
            f"Supported: {sorted(SUPPORTED_FORMATS)}"
        )
    return parser(path)


# ─── parsers ──────────────────────────────────────────────────────────────


def _parse_pdf(path: Path) -> list[PageContent]:
    return parse_pdf(path)


def _parse_txt(path: Path) -> list[PageContent]:
    text = path.read_text(encoding="utf-8", errors="replace")
    return [PageContent(
        page_number=1, markdown=text,
        has_text=bool(text.strip()), has_images=False, needs_ocr=False,
        width_pt=0.0, height_pt=0.0,
    )]


_H1_SPLIT_RE = re.compile(r"(?m)^#\s+")


def _parse_md(path: Path) -> list[PageContent]:
    text = path.read_text(encoding="utf-8", errors="replace")
    if not text.strip():
        return []
    # Split on top-level H1s — each section becomes a "page". If there are
    # no H1s, the whole file is one page.
    parts = _H1_SPLIT_RE.split(text)
    # Re-attach the '#' that the split consumed (skip the first which is
    # any preamble before the first H1).
    pages: list[str] = []
    if parts and parts[0].strip() and not text.lstrip().startswith("#"):
        pages.append(parts[0])
    pages.extend(["# " + p for p in parts[1:] if p.strip()])
    if not pages:
        pages = [text]
    return [
        PageContent(
            page_number=i + 1, markdown=p,
            has_text=True, has_images=False, needs_ocr=False,
            width_pt=0.0, height_pt=0.0,
        )
        for i, p in enumerate(pages)
    ]


def _parse_html(path: Path) -> list[PageContent]:
    from markdownify import markdownify
    raw = path.read_text(encoding="utf-8", errors="replace")
    md = markdownify(raw, heading_style="ATX")
    # Re-use the markdown splitter for H1-based pagination
    return _parse_md_string(md)


def _parse_md_string(text: str) -> list[PageContent]:
    if not text.strip():
        return []
    parts = _H1_SPLIT_RE.split(text)
    pages: list[str] = []
    if parts and parts[0].strip() and not text.lstrip().startswith("#"):
        pages.append(parts[0])
    pages.extend(["# " + p for p in parts[1:] if p.strip()])
    if not pages:
        pages = [text]
    return [
        PageContent(
            page_number=i + 1, markdown=p,
            has_text=True, has_images=False, needs_ocr=False,
            width_pt=0.0, height_pt=0.0,
        )
        for i, p in enumerate(pages)
    ]


def _parse_docx(path: Path) -> list[PageContent]:
    from docx import Document as DocxDocument  # type: ignore

    doc = DocxDocument(str(path))
    current_section: list[str] = []
    sections: list[list[str]] = []
    for para in doc.paragraphs:
        text = (para.text or "").rstrip()
        if not text:
            current_section.append("")
            continue
        style = (para.style.name or "").lower() if para.style else ""
        if style.startswith("heading 1"):
            # New section
            if current_section:
                sections.append(current_section)
            current_section = [f"# {text}"]
        elif style.startswith("heading 2"):
            current_section.append(f"## {text}")
        elif style.startswith("heading 3"):
            current_section.append(f"### {text}")
        elif style.startswith("heading"):
            current_section.append(f"#### {text}")
        else:
            current_section.append(text)
    if current_section:
        sections.append(current_section)
    if not sections:
        return []
    return [
        PageContent(
            page_number=i + 1,
            markdown="\n\n".join(line for line in sec if line is not None),
            has_text=True, has_images=False, needs_ocr=False,
            width_pt=0.0, height_pt=0.0,
        )
        for i, sec in enumerate(sections)
    ]


def _parse_pptx(path: Path) -> list[PageContent]:
    from pptx import Presentation  # type: ignore

    pres = Presentation(str(path))
    pages: list[PageContent] = []
    for i, slide in enumerate(pres.slides):
        lines: list[str] = []
        title = None
        if slide.shapes.title and slide.shapes.title.text:
            title = slide.shapes.title.text.strip()
            lines.append(f"# {title}")
        for shape in slide.shapes:
            if shape == slide.shapes.title:
                continue
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                run_text = "".join(r.text for r in para.runs).strip()
                if run_text:
                    bullet = "- " if para.level > 0 else ""
                    lines.append(bullet + run_text)
        md = "\n\n".join(lines)
        pages.append(PageContent(
            page_number=i + 1, markdown=md,
            has_text=bool(md.strip()), has_images=False,
            needs_ocr=not bool(md.strip()),
            width_pt=float(pres.slide_width or 0),
            height_pt=float(pres.slide_height or 0),
        ))
    return pages


def _parse_image(path: Path) -> list[PageContent]:
    """An image file is a single page with no extractable text — always
    routed through the OCR pipeline."""
    # Use PyMuPDF to get the pixel dimensions cleanly
    try:
        with pymupdf.open(str(path)) as doc:
            page = doc[0]
            width, height = float(page.rect.width), float(page.rect.height)
    except Exception:
        width = height = 0.0
    return [PageContent(
        page_number=1, markdown="",
        has_text=False, has_images=True, needs_ocr=True,
        width_pt=width, height_pt=height,
    )]


_PARSERS: dict[str, Callable[[Path], list[PageContent]]] = {
    "pdf": _parse_pdf,
    "txt": _parse_txt,
    "md": _parse_md,
    "markdown": _parse_md,
    "html": _parse_html,
    "htm": _parse_html,
    "docx": _parse_docx,
    "pptx": _parse_pptx,
}
for _ext in _IMAGE_EXTS:
    _PARSERS[_ext.lstrip(".")] = _parse_image


SUPPORTED_FORMATS = tuple(sorted(_PARSERS.keys()))


def render_page_image(path: str | Path, page_number: int, *, dpi: int = 200) -> bytes:
    """Render a page (or a whole image file) to PNG bytes. For non-PDF
    formats that PyMuPDF can open directly (images, EPUB, XPS) this works
    via fitz; for DOCX/PPTX/HTML/MD/TXT we don't render — those formats
    don't have a visual rendering and `needs_ocr` should never fire for
    them anyway."""
    p = Path(path)
    fmt = detect_format(p)
    if fmt == "pdf" or fmt in {"epub", "xps", "mobi", "cbz"} or f".{fmt}" in _IMAGE_EXTS:
        from kglite_docs.ingest.parser import render_page_png
        return render_page_png(p, page_number, dpi=dpi)
    raise ValueError(f"cannot render {fmt!r} to image — no visual layout")
